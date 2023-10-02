import json
import logging
import re
import requests
from PIL import Image
from PIL import ImageDraw
from pptx import Presentation
from pptx.util import Inches
import os
from dotenv import load_dotenv

# Load .env values into environment variables
load_dotenv()
# Get the API key from the environment variable
api_key = os.environ.get("API_KEY")


# Format the log message
logging.basicConfig(format='%(asctime)s %(levelname)s:%(name)s: %(message)s')
logger = logging.getLogger('dev')
logger.setLevel(logging.INFO)

text = {"outline": []}

headers = {
    "Content-Type": "application/json",
    "Authorization": f"Bearer {api_key}"
}

# filename = "text.pickle"
# infile = open(filename,'rb')
# text = pickle.load(infile)
# infile.close()
# logger.debug("Read text from file")



def initTextReq(input):
    logger.info("Requesting title, subtitle, and outline...")
    prompt = "Write a title, subtitle, and an outline for an essay about {input}."

    data = {
        "model": "text-davinci-003",
        "prompt": prompt.format(input=input),
        "temperature": 0,
        "max_tokens": 150
    }
    response = requests.post("https://api.openai.com/v1/completions", headers=headers, data=json.dumps(data))
    logger.debug("Response Text: " + response.json().get('choices')[0].get('text').strip())
    if response.status_code == 200:
        logger.info("Request successful.")
        return response.json().get('choices')[0].get('text')
    else:
        logger.error("Error: " + response.text)
        return response.json().get('choices')[0].get('text')


def quoteReq(input):
    logger.info("Requesting quote...")
    prompt = "Find a quote about {input}."
    data = {
        "model": "text-davinci-003",
        "prompt": prompt.format(input=input),
        "temperature": 0,
        "max_tokens": 40
    }
    response = requests.post("https://api.openai.com/v1/completions", headers=headers, data=json.dumps(data))
    logger.debug("Response Text: " + response.json().get('choices')[0].get('text').strip())
    if response.status_code == 200:
        logger.info("Request successful.")
        return response.json().get('choices')[0].get('text').strip()
    else:
        logger.error("Error: " + response.text)
        return response.json().get('choices')[0].get('text').strip()


def analogyReq(input):
    logger.info("Requesting analogy...")
    prompt = "Write an analogy about {input}."
    data = {
        "model": "text-davinci-003",
        "prompt": prompt.format(input=input),
        "temperature": 0,
        "max_tokens": 40
    }
    response = requests.post("https://api.openai.com/v1/completions", headers=headers, data=json.dumps(data))
    logger.debug("Response Text: " + response.json().get('choices')[0].get('text').strip())
    if response.status_code == 200:
        logger.info("Request successful.")
        return response.json().get('choices')[0].get('text').strip()
    else:
        logger.error("Error: " + response.text)
        return response.json().get('choices')[0].get('text').strip()


def bulletsReq(input, topic):
    logger.info("Requesting bullet points for " + topic + "...")
    prompt = "For a presentation about {input}, write some bullet points for a slide titled \"{topic}\""

    data = {
        "model": "text-davinci-003",
        "prompt": prompt.format(input=input, topic=topic),
        "temperature": 0,
        "max_tokens": 70
    }
    response = requests.post("https://api.openai.com/v1/completions", headers=headers, data=json.dumps(data))
    logger.debug("Response Text: " + str(response.json().get('choices')[0].get('text').strip()))
    if response.status_code == 200:
        logger.info("Request successful.")
        return response.json().get('choices')[0].get('text').strip()
    else:
        logger.error("Error: " + response.text)
        return response.json().get('choices')[0].get('text').strip()


def backImageReq(input):
    logger.info("Requesting image...")
    prompt = "Background image for a presentation slide about {input}"
    # Set the headers for the request
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }

    # Set the data for the request
    data = {
        "prompt": prompt.format(input=input),
        "n": 1,
        "size": "512x512"
    }

    # Make the request to the OpenAI API
    response = requests.post("https://api.openai.com/v1/images/generations", headers=headers, data=json.dumps(data))
    logger.debug("Response: " + str(response))
    # Check if the request was successful
    if response.status_code == 200:
        # Get the image data from the response
        image_data = response.json()["data"][0]["url"]
        # Download the image and save it to a file
        image = requests.get(image_data)
        open("image.jpg", "wb").write(image.content)
        logger.info("Image saved to file.")
        return response.status_code
    else:
        logger.error("Error: " + response.text)
        return response.status_code


def overlay_gradient():
    # Open the JPG image
    image = Image.open("image.jpg")

    # Create a new image with the same size as the original image
    gradient = Image.new("RGBA", image.size)

    # Create a linear gradient
    draw = ImageDraw.Draw(gradient)
    draw.linear_gradient((0, 0, image.width, image.height), (255, 0, 0), (0, 0, 255), "RGBA")

    # Overlay the gradient on the original image
    result = Image.alpha_composite(image, gradient)

    # Save the resulting image
    result.save("image.jpg")


def analyzeOutline(initText):
    global i, j
    title_match = re.search(r"Title: (.*?)\n", initText)
    subtitle_match = re.search(r"Subtitle: (.*?)\n", initText)
    outline_match = re.search(r"Outline:\n(.*)", initText, re.DOTALL)
    text["title"] = title_match.group(1) if title_match else None
    text["subtitle"] = subtitle_match.group(1) if subtitle_match else None
    outline = outline_match.group(1) if outline_match else None
    logger.debug("Title: " + text["title"])
    logger.debug("Subtitle: " + text["subtitle"])
    logger.debug("Outline: " + outline)
    text["quote"] = quoteReq(input)
    text["analogy"] = analogyReq(input)
    logger.debug("Original Outline: " + outline)
    # parse outline into parts
    outlineParts = re.split(r"[IVX]+\.", outline.lstrip())
    outlineParts[0] = outlineParts[0].lstrip()
    logger.debug("Outline Parts: " + str(outlineParts))
    i = -1
    for part in outlineParts:  # Section of outline (section)
        if part.strip():
            text["outline"].append({})  # create dict for section
            ++i
            lines = re.split(r"[A-Z]\.", part)
            logger.debug("Outline Lines: " + str(lines))
            k = 0  # counts number of topics added
            j = 0  # counts iterations of loop disregarding ''
            for line in lines:
                if line.strip():
                    if j == 0:  # First line (name of section)
                        text["outline"][i] = {"name": line.strip(), "topics": []}  # section
                        j = j + 1
                    else:  # Subsequent lines (topics)
                        text["outline"][i]["topics"].append({"name": line.strip(), "bullets": []})  # topic
                        bullets = bulletsReq(input, line.strip())
                        for l, bullet in enumerate(re.split("•", bullets)):
                            if bullet.strip():
                                if l < len(re.split("•", bullets)) - 1 or bullet.endswith('.'):
                                    text["outline"][i]["topics"][k]["bullets"].append(bullet.strip())
                        k = k + 1
                        j = j + 1
    logger.debug("Processed Outline: " + str(text["outline"]))





print("What should the powerpoint be about?")
input = input("Enter your prompt: ")

backImageReq(input)
# overlay_gradient()
analyzeOutline(initTextReq(input))



# # save text object to file
# text["input"] = input
# filename = "text.pickle"
# outfile = open(filename,'wb')
# pickle.dump(text,outfile)
# outfile.close()
# logger.info("Encoded text to file")



# create presentation

logger.info("Creating Presentation...")
pres = Presentation()

slide = pres.slides.add_slide(pres.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = text["title"]
subtitle.text = text["subtitle"]

left = top = Inches(0)
pic = slide.shapes.add_picture('image.jpg', left, top, width=pres.slide_width, height=pres.slide_height)

# This moves it to the background
slide.shapes._spTree.remove(pic._element)
slide.shapes._spTree.insert(2, pic._element)


logger.debug("Finished Title Slide")

# Analogy Slide
slide = pres.slides.add_slide(pres.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = text["analogy"]


logger.debug("Finished Analogy Slide")

# Iterate over the lists in the list

for i, section in enumerate(text["outline"]):
    # create title slide with name
    slide = pres.slides.add_slide(pres.slide_layouts[0])
    title = slide.shapes.title
    title.text = section["name"]
    for j, topic in enumerate(section["topics"]):
        logger.info(
            "Working on content slide " + str(j + 1) + "/" + str(len(section["topics"])) + " for topic " + str(i + 1) + "/" + str(
                len(text["outline"])))
        slide = pres.slides.add_slide(pres.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        title_shape.text = topic["name"]
        tf = body_shape.text_frame
        tf.text = "\n".join(topic["bullets"])


logger.info("Finished Content Slides")

# Quote Slide
slide = pres.slides.add_slide(pres.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = text["quote"]

logger.info("Finished Quote Slide")

# Save the presentation
pres.save('presentation.pptx')
logger.info("Finished Saving Presentation. Open presentation.pptx to view.")

# open the presentation
os.system("open presentation.pptx")